from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("deliverables")

NAVY = RGBColor(21, 24, 31)
INK = RGBColor(42, 45, 55)
MUTED = RGBColor(103, 111, 126)
LINE = RGBColor(220, 225, 233)
BLUE = RGBColor(77, 124, 251)
GREEN = RGBColor(27, 175, 122)
ORANGE = RGBColor(235, 104, 52)
YELLOW = RGBColor(237, 161, 0)
CYAN = RGBColor(8, 145, 178)
PINK = RGBColor(232, 123, 164)
RED = RGBColor(227, 73, 72)
LAVENDER = RGBColor(144, 133, 233)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(247, 249, 252)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=22, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.65, 0.45, 8.2, 0.55, 25, True, NAVY)
    if subtitle:
        add_text(slide, subtitle, 0.67, 1.02, 8.5, 0.35, 11, False, MUTED)


def add_footer(slide, label):
    add_text(slide, label, 0.65, 7.05, 4.2, 0.25, 8, False, MUTED)
    add_text(slide, "CEO Agent OS", 11.0, 7.05, 1.6, 0.25, 8, False, MUTED, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.adjustments[0] = 0.08
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, title, x + 0.22, y + 0.16, w - 0.35, 0.3, 13, True, NAVY)
    add_text(slide, body, x + 0.22, y + 0.56, w - 0.35, h - 0.7, 10, False, MUTED)


def add_bullets(slide, bullets, x, y, w, h, size=13, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Aptos"
        p.space_after = Pt(7)
    return box


def add_pill(slide, text, x, y, w, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    add_text(slide, text, x, y + 0.075, w, 0.2, 8.5, True, WHITE, PP_ALIGN.CENTER)


def add_flow(slide, items, x, y, w, color=BLUE):
    step_w = w / len(items)
    for i, item in enumerate(items):
        cx = x + i * step_w
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(y), Inches(step_w - 0.22), Inches(0.78))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LINE
        shape.adjustments[0] = 0.08
        add_text(slide, item, cx + 0.12, y + 0.2, step_w - 0.46, 0.26, 10, True, INK, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            add_text(slide, "->", cx + step_w - 0.18, y + 0.26, 0.25, 0.2, 11, True, color, PP_ALIGN.CENTER)


def cover(deck, title, subtitle, audience):
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    set_bg(slide, WHITE)
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.2), Inches(7.5))
    block.fill.solid()
    block.fill.fore_color.rgb = NAVY
    block.line.fill.background()
    add_text(slide, "CEO\nAgent\nOS", 0.65, 0.9, 2.8, 1.8, 31, True, WHITE)
    add_text(slide, audience, 0.68, 3.15, 2.7, 0.4, 12, False, RGBColor(184, 195, 214))
    add_text(slide, title, 4.75, 1.25, 7.2, 1.0, 29, True, NAVY)
    add_text(slide, subtitle, 4.78, 2.35, 6.9, 1.0, 15, False, MUTED)
    add_pill(slide, "Multi-agent operations", 4.8, 4.1, 1.85, BLUE)
    add_pill(slide, "Task execution", 6.85, 4.1, 1.45, GREEN)
    add_pill(slide, "Connected accounts", 8.5, 4.1, 1.65, ORANGE)
    add_footer(slide, audience)


def build_user_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(
        prs,
        "Product Walkthrough",
        "A simple operating dashboard where a CEO gives goals and specialized AI departments turn them into tasks, documents, drafts, and channel actions.",
        "User-facing deck",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "What This Project Is", "A local command center for delegating work to an automated company.")
    add_card(s, 0.75, 1.65, 3.75, 1.55, "One CEO inbox", "Users describe an outcome once. The CEO agent analyzes the ask, chooses the right department, and reports progress back in one run log.", BLUE)
    add_card(s, 4.85, 1.65, 3.75, 1.55, "Specialist departments", "Manager, HR, Developer, Analysis, Sales, Finance, SEO, Emails, and PR each have focused responsibilities and tools.", GREEN)
    add_card(s, 8.95, 1.65, 3.55, 1.55, "Saved deliverables", "Runs, documents, Linear tasks, costs, and transcripts remain available from the dashboard.", ORANGE)
    add_flow(s, ["Goal", "CEO decides", "Specialist acts", "Result saved"], 1.25, 4.6, 10.85)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Who Will Use It", "Designed for small teams that need a lightweight AI operations layer.")
    add_card(s, 0.8, 1.5, 3.6, 1.4, "Founder / CEO", "Delegates broad goals, checks outcomes, approves external actions, and monitors work across departments.", BLUE)
    add_card(s, 4.85, 1.5, 3.6, 1.4, "Operators", "Turn repeated workflows into requests, attach files, review documents, and archive completed work.", YELLOW)
    add_card(s, 8.9, 1.5, 3.6, 1.4, "Functional leads", "Use department pages directly when they already know the right specialist for a task.", GREEN)
    add_bullets(s, ["Best fit: solo founders, agencies, early-stage startup teams, and internal automation pilots.", "Not a replacement for final human approval on public posts, sent messages, hiring decisions, or financial commitments."], 1.05, 4.25, 11.1, 1.2, 15)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Core User Workflows", "The UI supports both broad delegation and direct specialist work.")
    add_card(s, 0.7, 1.35, 5.8, 1.05, "Overview run", "Submit a goal to the CEO agent. It decides who should act and coordinates specialist results.", BLUE)
    add_card(s, 6.85, 1.35, 5.8, 1.05, "Department run", "Open a department tab and ask that agent directly when the owner is already clear.", GREEN)
    add_card(s, 0.7, 2.7, 5.8, 1.05, "Attach context", "Upload PDF, DOCX, XLSX, CSV, markdown, JSON, logs, or code. The server extracts text and adds it to the prompt.", ORANGE)
    add_card(s, 6.85, 2.7, 5.8, 1.05, "Continue a run", "Reply inside a finished run to resume the same Claude session rather than starting from scratch.", LAVENDER)
    add_card(s, 0.7, 4.05, 5.8, 1.05, "Review outputs", "Open saved documents, inspect tool calls, view Linear tasks, and archive or delete runs.", CYAN)
    add_card(s, 6.85, 4.05, 5.8, 1.05, "Connect accounts", "Connect Gmail, Instagram, LinkedIn, Zernio, WhatsApp, image generation, and n8n workflows as needed.", RED)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Departments", "Each specialist has a narrow job and a bounded toolset.")
    departments = [
        ("Manager", "Breaks initiatives into Linear tasks", BLUE),
        ("HR", "Policies, onboarding, job descriptions", ORANGE),
        ("Developer", "Code and command execution in workspace", GREEN),
        ("Analysis", "Research and analytical reports", YELLOW),
        ("Sales", "Outreach, proposals, social drafts", PINK),
        ("Finance", "Budgets and financial write-ups", RGBColor(0, 131, 0)),
        ("SEO", "Keyword and competitor research", CYAN),
        ("Emails", "Inbox reading, drafts, sending when approved", LAVENDER),
        ("PR", "Announcements, pitches, public statements", RED),
    ]
    for i, (name, desc, color) in enumerate(departments):
        x = 0.75 + (i % 3) * 4.2
        y = 1.45 + (i // 3) * 1.45
        add_card(s, x, y, 3.75, 1.05, name, desc, color)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "What The Dashboard Shows", "The application is built around traceability.")
    add_bullets(s, [
        "Overview analytics: total runs, success rate, errors, documents, and Linear tasks.",
        "Run detail: status, cost, timestamps, streamed agent messages, tool calls, and tool results.",
        "Department pages: recent work for that department plus saved documents.",
        "Accounts page: connection status and setup guidance for supported integrations.",
        "Archive/delete controls: hide old tasks without losing history, or remove records permanently after confirmation.",
    ], 0.95, 1.55, 5.6, 4.2, 15)
    add_card(s, 7.0, 1.55, 2.5, 1.25, "Runs", "What happened and when", BLUE)
    add_card(s, 9.85, 1.55, 2.5, 1.25, "Costs", "Per-turn and total spend", YELLOW)
    add_card(s, 7.0, 3.15, 2.5, 1.25, "Documents", "Generated reports and drafts", GREEN)
    add_card(s, 9.85, 3.15, 2.5, 1.25, "Tasks", "Linear issues created", ORANGE)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Example Scenario", "From a vague business request to concrete work.")
    add_flow(s, ["User: launch campaign", "CEO delegates", "Sales drafts", "PR writes release", "Manager creates tasks"], 0.8, 1.65, 11.85)
    add_card(s, 1.0, 3.25, 3.5, 1.35, "Sales output", "Outreach sequence, proposal copy, social post draft, and needed customer context.", PINK)
    add_card(s, 4.9, 3.25, 3.5, 1.35, "PR output", "Press release, media pitch, LinkedIn-ready announcement draft.", RED)
    add_card(s, 8.8, 3.25, 3.5, 1.35, "Manager output", "Linear issues with acceptance criteria and ownership-ready descriptions.", BLUE)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "External Actions Need Approval", "The product is intentionally conservative around irreversible actions.")
    add_card(s, 0.9, 1.65, 3.5, 1.45, "Safe by default", "Email agents create drafts by default. Public posts and messages are described back before publishing.", GREEN)
    add_card(s, 4.9, 1.65, 3.5, 1.45, "Immediate actions", "Sending email, WhatsApp messages, Zernio DMs, and LinkedIn posts are treated as irreversible.", RED)
    add_card(s, 8.9, 1.65, 3.5, 1.45, "Configured accounts", "Tools appear only when the relevant token or API key is configured.", BLUE)
    add_bullets(s, ["Users keep control over public communication.", "The run log provides an audit trail for tool input and output.", "Unsupported or missing integrations are surfaced as setup states, not hidden failures."], 1.1, 4.1, 10.9, 1.4, 15)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Benefits", "Why this helps day-to-day operations.")
    add_bullets(s, [
        "Reduces context switching across planning, writing, research, email, social, and task tracking.",
        "Captures decisions and outputs in one persistent dashboard.",
        "Lets non-technical users route work by outcome instead of by tool.",
        "Makes agent work inspectable through logs, tool cards, documents, and task links.",
        "Keeps developer work sandboxed away from the app source by default.",
    ], 1.0, 1.5, 11.0, 4.3, 16)
    add_footer(s, "User-facing deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Roadmap Ideas", "Natural next steps for a production-grade version.")
    add_card(s, 0.85, 1.5, 3.7, 1.25, "Role-based access", "Separate admin, operator, and reviewer permissions.", BLUE)
    add_card(s, 4.85, 1.5, 3.7, 1.25, "Approval queues", "Formal review flow before sending, posting, or triggering external workflows.", GREEN)
    add_card(s, 8.85, 1.5, 3.7, 1.25, "More integrations", "CRM, calendar, analytics, storage, and ticketing connectors.", ORANGE)
    add_card(s, 0.85, 3.15, 3.7, 1.25, "Better observability", "Run metrics, latency, failure reasons, and cost budgets.", YELLOW)
    add_card(s, 4.85, 3.15, 3.7, 1.25, "Templates", "Reusable briefs for launches, hiring, campaigns, reports, and inbox triage.", PINK)
    add_card(s, 8.85, 3.15, 3.7, 1.25, "Deployment hardening", "Auth, encrypted tokens, job queues, and durable production storage.", RED)
    add_footer(s, "User-facing deck")

    OUT_DIR.mkdir(exist_ok=True)
    prs.save(OUT_DIR / "CEO-Agent-OS_User-Deck.pptx")


def build_tech_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(
        prs,
        "Technical Architecture",
        "A TypeScript/Node multi-agent orchestration app using Claude Agent SDK, Express, MCP tool servers, local persistence, and browser-based run monitoring.",
        "Developer deck",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "System At A Glance", "The app wraps a multi-agent runtime with a web dashboard and external tool connectors.")
    add_flow(s, ["Browser UI", "Express API", "CEO orchestrator", "Specialist agents", "MCP tools"], 0.85, 1.55, 11.7)
    add_card(s, 0.95, 3.0, 3.55, 1.35, "Frontend", "Static HTML/CSS/JS served from Express. Uses Lucide icons, marked, DOMPurify, GSAP, and Lenis from CDNs.", BLUE)
    add_card(s, 4.9, 3.0, 3.55, 1.35, "Backend", "Express server in TypeScript, run APIs, SSE streaming, upload parsing, account callbacks, and local stores.", GREEN)
    add_card(s, 8.85, 3.0, 3.55, 1.35, "Agent runtime", "Claude Agent SDK query(), custom Agent definitions, MCP servers, and explicit tool allowlists.", ORANGE)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Tech Stack", "Current dependencies and runtime choices from package.json.")
    add_card(s, 0.8, 1.35, 3.7, 1.2, "Language", "TypeScript 5.7, ESM modules, Node runtime, tsx for dev execution.", BLUE)
    add_card(s, 4.85, 1.35, 3.7, 1.2, "Server", "Express 4, multer, native fetch, Server-Sent Events for live run updates.", GREEN)
    add_card(s, 8.9, 1.35, 3.7, 1.2, "Agents", "@anthropic-ai/claude-agent-sdk with custom subagents and MCP servers.", ORANGE)
    add_card(s, 0.8, 2.9, 3.7, 1.2, "Data parsing", "exceljs, pdf-parse, mammoth, plain text/code formats.", YELLOW)
    add_card(s, 4.85, 2.9, 3.7, 1.2, "Integrations", "Google APIs, Linear GraphQL, Meta Graph APIs, LinkedIn REST, Zernio API, n8n webhooks.", CYAN)
    add_card(s, 8.9, 2.9, 3.7, 1.2, "Validation", "zod schemas for MCP tool inputs.", LAVENDER)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Agent Architecture", "CEO delegates; specialists execute with narrow prompts and tools.")
    add_bullets(s, [
        "DEPARTMENTS defines the product-facing department roster, labels, icons, taglines, and identity colors.",
        "buildAgentsRegistry() creates Manager, HR, Developer, Analysis, Sales, Finance, SEO, Emails, and PR agents.",
        "CEO_SYSTEM_PROMPT instructs the top-level agent to analyze, delegate, and summarize instead of doing specialist work itself.",
        "Specialists use background: false to reduce lost-result behavior in delegated tool calls.",
        "runSpecialistAgent() bypasses the CEO for department-page direct runs.",
    ], 0.9, 1.45, 11.5, 4.8, 14)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Runtime Flow", "How a user request turns into saved state.")
    add_flow(s, ["POST /api/runs", "createRun()", "query()", "drainQuery()", "appendEvent()", "finishRun()"], 0.65, 1.45, 12.1)
    add_card(s, 0.95, 3.05, 3.6, 1.45, "Streaming", "SDK messages are converted into text, tool_use, tool_result, and done events, then streamed to the browser by SSE.", BLUE)
    add_card(s, 4.9, 3.05, 3.6, 1.45, "Task capture", "Linear task refs are parsed from tool results and stored on the run record.", GREEN)
    add_card(s, 8.85, 3.05, 3.6, 1.45, "Resume", "Captured SDK session IDs let /api/runs/:id/reply continue a finished run in place.", ORANGE)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "API Surface", "Express routes are small and resource-oriented.")
    add_bullets(s, [
        "POST /api/runs: start a CEO-routed run.",
        "POST /api/agents/:key/runs: start a direct specialist run.",
        "POST /api/runs/:id/reply: continue a captured SDK session.",
        "GET /api/runs and GET /api/runs/:id: list and inspect persisted run records.",
        "GET /api/runs/:id/stream: stream run updates via text/event-stream.",
        "GET /api/departments, /api/analytics, /api/documents, /api/accounts: dashboard data.",
        "POST /api/uploads: parse supported attachment files into text context.",
        "POST /api/automation/*: API-key-gated run creation for external automations.",
    ], 0.95, 1.35, 11.4, 5.2, 13)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "MCP Tool Servers", "Each external capability is exposed through bounded tools.")
    add_card(s, 0.75, 1.35, 3.9, 1.2, "linear", "create_linear_task, list_linear_tasks via Linear GraphQL.", BLUE)
    add_card(s, 4.75, 1.35, 3.9, 1.2, "documents", "create_{department}_document tools persisted under workspace/.documents.", GREEN)
    add_card(s, 8.75, 1.35, 3.9, 1.2, "gmail", "list/read email, create draft, send when explicitly instructed.", ORANGE)
    add_card(s, 0.75, 2.9, 3.9, 1.2, "social", "Instagram DMs, LinkedIn organization posts, Zernio posts and conversations.", PINK)
    add_card(s, 4.75, 2.9, 3.9, 1.2, "whatsapp", "Template and freeform sends through WhatsApp Cloud API.", CYAN)
    add_card(s, 8.75, 2.9, 3.9, 1.2, "automation/media", "n8n allowlisted workflows, OpenAI image generation, local post images.", LAVENDER)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Data And Persistence", "Simple local files back the single-user dashboard.")
    add_bullets(s, [
        "data/runs.json stores run records, events, status, costs, task refs, archive flags, and SDK session IDs.",
        "workspace/.documents/index.json stores generated markdown deliverables.",
        "data/*-token.json stores OAuth connection tokens for Gmail, Instagram, and LinkedIn.",
        "Workspace isolation is intentional: agent cwd is workspace/, not the app source tree.",
        "Document writes live under workspace because delegated subagent writes outside cwd were observed to be unreliable.",
    ], 0.9, 1.45, 11.4, 4.8, 14)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Frontend Design", "A static app that renders from API state.")
    add_card(s, 0.8, 1.4, 3.7, 1.25, "State model", "One state object tracks view, theme, departments, accounts, runs, selected run, docs, and live tasks.", BLUE)
    add_card(s, 4.85, 1.4, 3.7, 1.25, "Views", "Overview dashboard, department pages, run detail, accounts page, document popups.", GREEN)
    add_card(s, 8.9, 1.4, 3.7, 1.25, "Rendering", "Vanilla JS string templates with event delegation after each render.", ORANGE)
    add_card(s, 0.8, 3.05, 3.7, 1.25, "Safety", "Agent markdown is rendered with marked and sanitized with DOMPurify.", RED)
    add_card(s, 4.85, 3.05, 3.7, 1.25, "Live UX", "SSE updates, tool cards, auto-follow log behavior, reply form, and run actions.", CYAN)
    add_card(s, 8.9, 3.05, 3.7, 1.25, "Responsive UI", "Sidebar drawer on mobile, fixed shell layout on desktop, light/dark theme support.", LAVENDER)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Configuration", "Runtime capability is controlled by environment variables and token files.")
    add_bullets(s, [
        "Core: PORT, WEBHOOK_URL, AUTOMATION_API_KEY.",
        "Linear: LINEAR_API_KEY, LINEAR_TEAM_ID.",
        "Google/Gmail: GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET, GOOGLE_REDIRECT_URI.",
        "Meta/Instagram/WhatsApp: META_APP_ID, META_APP_SECRET, META_REDIRECT_URI, WHATSAPP_*.",
        "LinkedIn: LINKEDIN_CLIENT_ID, LINKEDIN_CLIENT_SECRET, LINKEDIN_REDIRECT_URI.",
        "Zernio and image posting: ZERNIO_API_KEY, OPENAI_API_KEY.",
        "n8n: N8N_WORKFLOWS as an allowlist mapping workflow names to webhook URLs.",
    ], 0.95, 1.45, 11.5, 4.9, 14)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Security And Reliability Notes", "Important implementation decisions developers should preserve.")
    add_bullets(s, [
        "Developer agent runs inside workspace/ so Bash and file tools do not touch the app source tree.",
        "n8n trigger accepts only admin-configured workflow names, preventing arbitrary URL calls.",
        "External sends/posts are described as irreversible and prompts default to drafts or approval-first behavior.",
        "The automation API is gated by X-API-Key and separate from browser run creation.",
        "Uploaded files are parsed in memory and discarded; text is appended to the agent prompt.",
        "Known reliability mitigation: SDK task_notification events are handled because backgrounded subagent outcomes can arrive there.",
    ], 0.95, 1.35, 11.4, 5.1, 13.5)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Build, Run, Verify", "Commands available in the current project.")
    add_card(s, 0.9, 1.5, 3.65, 1.2, "Install", "npm install", BLUE)
    add_card(s, 4.85, 1.5, 3.65, 1.2, "Run web app", "npm start\nServes http://localhost:3000 by default.", GREEN)
    add_card(s, 8.8, 1.5, 3.65, 1.2, "Run CLI", "npm run cli -- \"<goal>\"", ORANGE)
    add_card(s, 0.9, 3.15, 3.65, 1.2, "Compile", "npm run build", CYAN)
    add_card(s, 4.85, 3.15, 3.65, 1.2, "Typecheck", "npm run typecheck", LAVENDER)
    add_card(s, 8.8, 3.15, 3.65, 1.2, "Open risks", "No automated test suite is defined yet; add focused tests around stores, parsers, and route behavior.", RED)
    add_footer(s, "Developer deck")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Next Engineering Steps", "Work that would make the system production-ready.")
    add_bullets(s, [
        "Move from local JSON/token files to durable storage with encryption for credentials.",
        "Add authentication, authorization, and per-action approval records.",
        "Add a job queue so long agent runs survive process restarts.",
        "Create automated tests for store persistence, attachment parsing, route contracts, and integration error handling.",
        "Add observability for run latency, tool failures, token/cost budgets, and integration health.",
        "Modularize frontend rendering if the UI continues to grow beyond the current static app pattern.",
    ], 0.95, 1.5, 11.4, 4.8, 15)
    add_footer(s, "Developer deck")

    OUT_DIR.mkdir(exist_ok=True)
    prs.save(OUT_DIR / "CEO-Agent-OS_Technical-Deck.pptx")


if __name__ == "__main__":
    build_user_deck()
    build_tech_deck()
    print("Created:")
    print(OUT_DIR / "CEO-Agent-OS_User-Deck.pptx")
    print(OUT_DIR / "CEO-Agent-OS_Technical-Deck.pptx")
